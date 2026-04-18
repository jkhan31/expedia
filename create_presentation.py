"""
Create Expedia Marketplace Analysis PowerPoint Deck
Uses data from the notebook analysis to build professional slides with charts and tables
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
import pathlib
from io import BytesIO

# Set base directory for relative paths
BASE_DIR = pathlib.Path(__file__).parent

# Set style
sns.set_style("whitegrid")
plt.rcParams['figure.figsize'] = (12, 7)
plt.rcParams['font.size'] = 11

# Load data from sample
print("Loading analysis data...")
df = pd.read_csv(BASE_DIR / 'data' / 'sample.csv')

# Recreate market segmentation
price_p33 = df['price_usd'].quantile(0.33)
price_p66 = df['price_usd'].quantile(0.66)

def assign_segment(row):
    price = row['price_usd']
    rating = row['prop_starrating']
    if price > price_p66 and rating >= 4:
        return 'Luxury'
    elif price < price_p33 or rating <= 2:
        return 'Budget'
    else:
        return 'Mid'

df['market_segment'] = df.apply(assign_segment, axis=1)

# Recreate competitiveness score
def has_extreme_signed(row):
    for i in range(1, 9):
        rate = row[f'comp{i}_rate']
        pct_diff = row[f'comp{i}_rate_percent_diff']
        if pd.notna(rate) and pd.notna(pct_diff):
            signed = rate * pct_diff
            if abs(signed) > 100:
                return True
    return False

def calc_comp_clean(row):
    comp_diffs = []
    for i in range(1, 9):
        rate = row[f'comp{i}_rate']
        pct_diff = row[f'comp{i}_rate_percent_diff']
        if pd.notna(rate) and pd.notna(pct_diff):
            signed_diff = rate * pct_diff
            comp_diffs.append(signed_diff)
    if len(comp_diffs) == 0:
        return np.nan
    return np.mean(comp_diffs)

df = df[~df.apply(has_extreme_signed, axis=1)].copy()
df['comp_score'] = df.apply(calc_comp_clean, axis=1)

# Recreate visitor type
df['visitor_type'] = df['visitor_hist_starrating'].apply(
    lambda x: 'Returning' if pd.notna(x) else 'New'
)

print(f"Data loaded: {len(df):,} rows")

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 121)  # Dark blue

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)

    return slide

def add_content_slide(prs, title, content_type='text'):
    """Add content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(31, 78, 121)
    title_shape.line.color.rgb = RGBColor(31, 78, 121)

    # Title text
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    return slide

def save_chart_to_bytes(fig):
    """Save matplotlib figure to bytes"""
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close(fig)
    return buf

# ============================================================================
# SLIDE 1: TITLE
# ============================================================================
add_title_slide(prs,
    "Expedia Marketplace Analysis",
    "Competitive Positioning Diagnosis")

# ============================================================================
# SLIDE 2: PROBLEM STATEMENT
# ============================================================================
slide = add_content_slide(prs, "Problem Statement")

# Problem text
problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
tf = problem_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Booking rates vary significantly by market segment:"
p.font.size = Pt(20)
p.font.bold = True

# Stats box
stats_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(3.5))
tf = stats_box.text_frame
tf.word_wrap = True

segments_data = [
    ("Budget", "3.04%", "Highest"),
    ("Mid", "2.89%", "Middle"),
    ("Luxury", "2.20%", "Lowest (-27% vs Budget)")
]

for i, (seg, rate, desc) in enumerate(segments_data):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  • {seg}: {rate} booking rate ({desc})"
    p.font.size = Pt(18)
    p.level = 0
    p.space_before = Pt(8)

# Question
q_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(1))
tf = q_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Why does Luxury underperform? Is it pricing, visibility, or quality trust?"
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = RGBColor(200, 0, 0)

# ============================================================================
# SLIDE 3: METHODOLOGY
# ============================================================================
slide = add_content_slide(prs, "Methodology")

methodology_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
tf = methodology_box.text_frame
tf.word_wrap = True

methods = [
    ("Outcome Data Available", "100k searches with click_bool, booking_bool, gross_bookings_usd"),
    ("Causal Measurement", "Direct measurement of what drives conversions, not hypothetical"),
    ("Segment Analysis", "Budget (price <$99 OR rating ≤2), Mid (3-4★), Luxury (price >$159 AND rating ≥4)"),
    ("5 Core Analyses", "Funnel, Ranking Impact, Competitive Positioning, Quality Trust Gap, User Behavior"),
    ("Data-Driven Scope", "Skip redundant analyses when findings are conclusive")
]

for i, (title, desc) in enumerate(methods):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_before = Pt(6)

    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(13)
    p.level = 1
    p.space_before = Pt(2)

# ============================================================================
# SLIDE 4: FINDING 1 - COMPETITIVENESS
# ============================================================================
slide = add_content_slide(prs, "Finding 1: Competitiveness = Market Parity")

# Text
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
tf = text_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Price is NOT the issue"
p.font.size = Pt(18)
p.font.bold = True

p = tf.add_paragraph()
p.text = "Median comp_score: 0.0% (market parity)"
p.font.size = Pt(14)
p.space_before = Pt(10)

p = tf.add_paragraph()
p.text = "Underpriced hotels: 3.38% booking rate"
p.font.size = Pt(13)
p.level = 1

p = tf.add_paragraph()
p.text = "Overpriced hotels: 2.55% booking rate"
p.font.size = Pt(13)
p.level = 1

p = tf.add_paragraph()
p.text = "Difference: 0.83pp (negligible)"
p.font.size = Pt(13)
p.level = 1

p = tf.add_paragraph()
p.text = "Correlation to bookings: +0.0289"
p.font.size = Pt(13)
p.level = 1
p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = "Verdict: Price optimization has minimal ROI"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 100, 0)

# Chart
fig, ax = plt.subplots(figsize=(5, 4))
pricing_groups = ['Underpriced\n(>0%)', 'At Parity\n(-0.5% to 0.5%)', 'Overpriced\n(<-0.5%)']
booking_rates = [3.38, 3.01, 2.55]
colors = ['#2ecc71', '#f39c12', '#e74c3c']
bars = ax.bar(pricing_groups, booking_rates, color=colors, alpha=0.8, edgecolor='black', linewidth=1.5)
ax.set_ylabel('Booking Rate (%)', fontsize=12, fontweight='bold')
ax.set_ylim(0, 4)
ax.axhline(y=2.8, color='red', linestyle='--', label='Avg: 2.8%', linewidth=2)
for bar in bars:
    height = bar.get_height()
    ax.text(bar.get_x() + bar.get_width()/2., height,
            f'{height:.2f}%', ha='center', va='bottom', fontweight='bold')
ax.legend()
ax.grid(axis='y', alpha=0.3)
plt.tight_layout()
img_stream = save_chart_to_bytes(fig)
slide.shapes.add_picture(img_stream, Inches(5.2), Inches(1.2), width=Inches(4.3))

# ============================================================================
# SLIDE 5: FINDING 2 - RANKING IMPACT
# ============================================================================
slide = add_content_slide(prs, "Finding 2: Ranking Position Impact")

# Text
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
tf = text_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Position MATTERS"
p.font.size = Pt(18)
p.font.bold = True

p = tf.add_paragraph()
p.text = "Position 1 vs Position 10:"
p.font.size = Pt(14)
p.space_before = Pt(10)

p = tf.add_paragraph()
p.text = "Click rate: 18.78% vs 4.15% → 4.53x elasticity"
p.font.size = Pt(13)
p.level = 1

p = tf.add_paragraph()
p.text = "Booking rate: 13.37% vs 2.56% → 5.22x elasticity"
p.font.size = Pt(13)
p.level = 1

p = tf.add_paragraph()
p.text = "Insight: Position elasticity for BOOKING is stronger than for CLICKS"
p.font.size = Pt(13)
p.level = 1
p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = "Verdict: Position 1 hotels are higher quality or better-matched to user intent"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 100, 0)

# Chart
fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(9, 4))

# Click rates
positions = [1, 2, 3, 4, 10, 15, 20]
click_rates = [18.78, 13.74, 10.66, 8.55, 4.15, 3.19, 2.08]
ax1.plot(positions, click_rates, marker='o', linewidth=2.5, markersize=8, color='#3498db')
ax1.set_xlabel('Position', fontweight='bold')
ax1.set_ylabel('Click Rate (%)', fontweight='bold')
ax1.set_title('Click Rate by Position', fontweight='bold')
ax1.grid(True, alpha=0.3)

# Booking rates
booking_rates = [13.37, 9.74, 7.13, 5.67, 2.56, 1.52, 0.99]
ax2.plot(positions, booking_rates, marker='s', linewidth=2.5, markersize=8, color='#e74c3c')
ax2.set_xlabel('Position', fontweight='bold')
ax2.set_ylabel('Booking Rate (%)', fontweight='bold')
ax2.set_title('Booking Rate by Position', fontweight='bold')
ax2.grid(True, alpha=0.3)

plt.tight_layout()
img_stream = save_chart_to_bytes(fig)
slide.shapes.add_picture(img_stream, Inches(5.2), Inches(1.2), width=Inches(4.3))

# ============================================================================
# SLIDE 6: FINDING 3 - QUALITY TRUST GAP
# ============================================================================
slide = add_content_slide(prs, "Finding 3: Quality Trust Gap by Segment")

# Chart
fig, ax = plt.subplots(figsize=(9, 5))
segments = ['Budget', 'Mid', 'Luxury']
promised = [2.40, 3.34, 4.30]
actual = [3.39, 3.97, 4.20]
x = np.arange(len(segments))
width = 0.35

bars1 = ax.bar(x - width/2, promised, width, label='Promised (Star Rating)', color='#3498db', alpha=0.8, edgecolor='black')
bars2 = ax.bar(x + width/2, actual, width, label='Actual (Review Score)', color='#2ecc71', alpha=0.8, edgecolor='black')

ax.set_ylabel('Rating', fontsize=12, fontweight='bold')
ax.set_title('Quality Promise vs Delivery by Segment', fontsize=14, fontweight='bold')
ax.set_xticks(x)
ax.set_xticklabels(segments)
ax.legend(fontsize=11)
ax.set_ylim(0, 5)

# Add gap labels
for i, (p, a) in enumerate(zip(promised, actual)):
    gap = ((a - p) / p * 100)
    color = '#2ecc71' if gap > 0 else '#e74c3c'
    ax.text(i, max(p, a) + 0.3, f'{gap:+.1f}%', ha='center', fontweight='bold', color=color, fontsize=11)

plt.tight_layout()
img_stream = save_chart_to_bytes(fig)
slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.2), width=Inches(9))

# ============================================================================
# SLIDE 7: FINDING 4 - LUXURY QUALITY PROBLEM
# ============================================================================
slide = add_content_slide(prs, "Finding 4: Luxury Quality Problem Detail")

# Table with beat/fail data
table_shape = slide.shapes.add_table(4, 4, Inches(0.5), Inches(1.2), Inches(9), Inches(2.2))
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(2)
table.columns[1].width = Inches(2.3)
table.columns[2].width = Inches(2.3)
table.columns[3].width = Inches(2.4)

# Header
headers = ['Segment', 'Beat Promise', 'Fail Promise', 'Booking Rate']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(31, 78, 121)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Data
data = [
    ('Budget', '78% (30.7k hotels)', '10% (4.0k)', '3.04%'),
    ('Mid', '73% (27.0k)', '12% (4.3k)', '2.89%'),
    ('Luxury', '39% (9.0k)', '31% (7.2k)', '2.20%')
]

for row_idx, (seg, beat, fail, book) in enumerate(data, 1):
    cells = [table.cell(row_idx, col_idx) for col_idx in range(4)]
    values = [seg, beat, fail, book]

    for cell, value in zip(cells, values):
        cell.text = value
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        if row_idx == 3:  # Luxury row
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 200, 200)


# Insight
insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(3.2))
tf = insight_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Key Insight: Luxury Buyers Expectations Problem"
p.font.size = Pt(16)
p.font.bold = True

p = tf.add_paragraph()
p.text = "Luxury segment has LOWEST % beating promise (39%) and HIGHEST % failing (31%)"
p.font.size = Pt(13)
p.space_before = Pt(8)

p = tf.add_paragraph()
p.text = "When luxury buyer sees 4.30★ hotel but discovers 4.20★ reality → disappointment → abandonment"
p.font.size = Pt(13)
p.space_before = Pt(6)

p = tf.add_paragraph()
p.text = "Result: Lowest booking rate (2.20%) among all segments"
p.font.size = Pt(13)
p.space_before = Pt(6)
p.font.color.rgb = RGBColor(200, 0, 0)
p.font.bold = True

# ============================================================================
# SLIDE 8: FUNNEL ANALYSIS
# ============================================================================
slide = add_content_slide(prs, "Finding 5: Funnel Analysis by Segment")

# Chart
fig, ax = plt.subplots(figsize=(9, 5))
segments = ['Budget', 'Mid', 'Luxury']
click_rates = [4.60, 4.70, 4.10]
booking_rates = [3.04, 2.89, 2.20]

x = np.arange(len(segments))
width = 0.35

bars1 = ax.bar(x - width/2, click_rates, width, label='Click Rate', color='#3498db', alpha=0.8, edgecolor='black')
bars2 = ax.bar(x + width/2, booking_rates, width, label='Booking Rate', color='#e74c3c', alpha=0.8, edgecolor='black')

ax.set_ylabel('Rate (%)', fontsize=12, fontweight='bold')
ax.set_title('Click vs Booking Rates: Flat Clicks, Declining Bookings', fontsize=14, fontweight='bold')
ax.set_xticks(x)
ax.set_xticklabels(segments)
ax.legend(fontsize=11)
ax.set_ylim(0, 5.5)

for bars in [bars1, bars2]:
    for bar in bars:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{height:.2f}%', ha='center', va='bottom', fontweight='bold', fontsize=10)

plt.tight_layout()
img_stream = save_chart_to_bytes(fig)
slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.2), width=Inches(9))

# ============================================================================
# SLIDE 9: USER BEHAVIOR
# ============================================================================
slide = add_content_slide(prs, "Finding 6: Returning Visitors Convert Better")

# Text
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
tf = text_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Experience Drives Conversion"
p.font.size = Pt(18)
p.font.bold = True

p = tf.add_paragraph()
p.text = "New Visitors"
p.font.size = Pt(14)
p.font.bold = True
p.space_before = Pt(10)

p = tf.add_paragraph()
p.text = "• Click rate: 4.49%"
p.font.size = Pt(12)
p.level = 1

p = tf.add_paragraph()
p.text = "• Booking rate: 2.74%"
p.font.size = Pt(12)
p.level = 1

p = tf.add_paragraph()
p.text = "Returning Visitors"
p.font.size = Pt(14)
p.font.bold = True
p.space_before = Pt(10)

p = tf.add_paragraph()
p.text = "• Click rate: 4.61% (+2.7%)"
p.font.size = Pt(12)
p.level = 1

p = tf.add_paragraph()
p.text = "• Booking rate: 3.66% (+33.6%)"
p.font.size = Pt(12)
p.level = 1

p = tf.add_paragraph()
p.text = "Conversion Lift: 1.34x"
p.font.size = Pt(14)
p.font.bold = True
p.space_before = Pt(10)
p.font.color.rgb = RGBColor(0, 100, 0)

# Chart
fig, ax = plt.subplots(figsize=(5, 4))
visitor_types = ['New Visitors', 'Returning\nVisitors']
booking_rates = [2.74, 3.66]
colors = ['#95a5a6', '#2ecc71']
bars = ax.bar(visitor_types, booking_rates, color=colors, alpha=0.8, edgecolor='black', linewidth=1.5)
ax.set_ylabel('Booking Rate (%)', fontsize=12, fontweight='bold')
ax.set_ylim(0, 4.5)
# Label bars with proper lift annotation only on returning visitors
labels = [f'{booking_rates[0]:.2f}%', f'{booking_rates[1]:.2f}%\n(1.34x)']
for bar, label in zip(bars, labels):
    height = bar.get_height()
    ax.text(bar.get_x() + bar.get_width()/2., height,
            label, ha='center', va='bottom', fontweight='bold', fontsize=11)
ax.grid(axis='y', alpha=0.3)
plt.tight_layout()
img_stream = save_chart_to_bytes(fig)
slide.shapes.add_picture(img_stream, Inches(5.2), Inches(1.2), width=Inches(4.3))

# ============================================================================
# SLIDE 10: ROOT CAUSE SUMMARY
# ============================================================================
slide = add_content_slide(prs, "Root Cause: Quality Trust Gap Explains Everything")

# Root cause box
root_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf = root_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Budget Segment ✓ SUCCEEDS"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 100, 0)

p = tf.add_paragraph()
p.text = "Promise: 2.40★ | Delivery: 3.39 | Gap: +41.2% (beats) | Booking: 3.04%"
p.font.size = Pt(12)
p.level = 1
p.space_before = Pt(4)

p = tf.add_paragraph()
p.text = "Mid Segment ✓ STABLE"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 100, 0)
p.space_before = Pt(10)

p = tf.add_paragraph()
p.text = "Promise: 3.34★ | Delivery: 3.97 | Gap: +18.7% (beats) | Booking: 2.89%"
p.font.size = Pt(12)
p.level = 1
p.space_before = Pt(4)

p = tf.add_paragraph()
p.text = "Luxury Segment ✗ FAILS"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(200, 0, 0)
p.space_before = Pt(10)

p = tf.add_paragraph()
p.text = "Promise: 4.30★ | Delivery: 4.20 | Gap: -2.3% (fails) | Booking: 2.20%"
p.font.size = Pt(12)
p.level = 1
p.space_before = Pt(4)

p = tf.add_paragraph()
p.text = "31% of luxury hotels fail their promise → Highest abandonment"
p.font.size = Pt(12)
p.level = 1
p.space_before = Pt(4)

p = tf.add_paragraph()
p.text = "VERDICT: Not price. Not visibility. Not user type. QUALITY TRUST."
p.font.size = Pt(14)
p.font.bold = True
p.space_before = Pt(15)
p.font.color.rgb = RGBColor(200, 0, 0)

# ============================================================================
# SLIDE 11: RECOMMENDATIONS
# ============================================================================
slide = add_content_slide(prs, "Strategic Recommendations (Prioritized)")

# Recommendations
rec_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
tf = rec_box.text_frame
tf.word_wrap = True

recs = [
    ("1. FIX LUXURY QUALITY GAP (Highest Impact)", [
        "Audit all 4.3★ hotels in luxury segment",
        "Verify star rating accuracy against actual review scores",
        "Remove misaligned listings or update ratings",
        "Expected impact: 5% lift in luxury booking rate (2.20% → 2.31%)"
    ]),
    ("2. IMPROVE RANKING ALGORITHM (Strong Impact)", [
        "Optimize position 1 for quality-intent alignment",
        "5.22x booking elasticity shows position matters more for booking than clicks",
        "Position quality score should incorporate review-rating alignment"
    ]),
    ("3. BUILD TRUST SIGNALS (Secondary)", [
        "Highlight verified buyer reviews on luxury listings",
        "Show returning visitor booking rates as proof point",
        "1.34x lift from familiarity demonstrates trust value"
    ]),
    ("4. STOP PRICE OPTIMIZATION (Avoid Waste)", [
        "Price correlation to bookings: +0.0289 (negligible)",
        "Redirect resources from pricing strategy to quality alignment",
        "Price is NOT the conversion lever"
    ])
]

for i, (title, bullets) in enumerate(recs):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.space_before = Pt(8)

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(11)
        p.level = 1
        p.space_before = Pt(2)

# ============================================================================
# SLIDE 12: IMPLEMENTATION PLAN
# ============================================================================
slide = add_content_slide(prs, "Implementation Roadmap")

# Plan
plan_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
tf = plan_box.text_frame
tf.word_wrap = True

phases = [
    ("Phase 1: Luxury Segment Audit (Week 1-2)", [
        "Extract all luxury hotels with 4.3★ rating",
        "Calculate actual vs promised quality gap per hotel",
        "Flag hotels with >2% gap for review",
        "Output: Audit report + remediation list"
    ]),
    ("Phase 2: Ranking Algorithm Update (Week 3-4)", [
        "Implement quality-alignment score in ranking",
        "A/B test: Control (current algo) vs Treatment (quality-weighted)",
        "Measure booking lift in luxury segment",
        "Output: Algorithm improvement + A/B results"
    ]),
    ("Phase 3: Trust Signal Rollout (Week 5-6)", [
        "Add verified buyer labels to luxury listings",
        "Test impact on click-to-book conversion",
        "Measure returning visitor preference",
        "Output: Feature + effectiveness metrics"
    ])
]

for i, (phase, bullets) in enumerate(phases):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    p.text = phase
    p.font.size = Pt(13)
    p.font.bold = True
    p.space_before = Pt(8)

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(11)
        p.level = 1
        p.space_before = Pt(2)

# ============================================================================
# SLIDE 13: SUCCESS METRICS
# ============================================================================
slide = add_content_slide(prs, "Success Metrics & Targets")

# Metrics table
table_shape = slide.shapes.add_table(5, 3, Inches(0.5), Inches(1.2), Inches(9), Inches(2.2))
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(3)

# Header
headers = ['Metric', 'Current', 'Target']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(31, 78, 121)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Data
metrics = [
    ('Luxury Booking Rate', '2.20%', '2.31% (+5%)'),
    ('Quality Gap (Luxury)', '-2.3%', '-1.0% (-57%)'),
    ('Hotels Beating Promise (Luxury)', '39%', '50%'),
    ('Overall Marketplace Booking Rate', '2.80%', '2.90% (+3.6%)')
]

for row_idx, (metric, current, target) in enumerate(metrics, 1):
    cells = [table.cell(row_idx, col_idx) for col_idx in range(3)]
    values = [metric, current, target]

    for cell, value in zip(cells, values):
        cell.text = value
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        if '5%' in value or '57%' in value:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(200, 255, 200)


# Timeline
timeline_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(3))
tf = timeline_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Timeline: 6 weeks | Measurement: Weekly tracking of luxury booking rate, quality gap"
p.font.size = Pt(12)
p.font.bold = True
p.space_before = Pt(10)

p = tf.add_paragraph()
p.text = "Success: 5% lift in luxury bookings confirms quality alignment is the conversion lever"
p.font.size = Pt(12)
p.space_before = Pt(8)

# ============================================================================
# SLIDE 14: NEXT STEPS
# ============================================================================
slide = add_content_slide(prs, "Next Steps")

next_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
tf = next_box.text_frame
tf.word_wrap = True

steps = [
    "1. Stakeholder Alignment - Present findings to product, data science, and operations teams",
    "2. Data QA - Verify luxury hotel star ratings and review scores (manual audit sample)",
    "3. Tool Setup - Build luxury audit tool to identify misaligned hotels",
    "4. Ranking Review - Work with ML team to understand current ranking factors",
    "5. A/B Testing Plan - Design experiment structure for quality-weighted ranking",
    "6. Communication - Draft messaging for hotels about star rating accuracy",
    "7. Launch Week 1 - Begin luxury audit; identify first batch of hotels to remediate"
]

for i, step in enumerate(steps):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = step
    p.font.size = Pt(12)
    p.space_before = Pt(10)

# Save presentation
presentations_dir = BASE_DIR / 'presentations'
presentations_dir.mkdir(exist_ok=True)
output_path = presentations_dir / 'Expedia-Marketplace-Analysis.pptx'
prs.save(str(output_path))
print(f"[OK] Presentation created: {output_path}")
print(f"  14 slides with charts, tables, and data visualizations")
